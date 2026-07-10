"""FoxSay 文档解析路由器。

根据文件类型自动分发到对应的解析器：
- .pdf → 双轨制：电子版 → Docling，扫描件 → MinerU V4 云端
- .docx/.xlsx/.html → MarkItDown 轻量解析
- .ppt/.pptx → python-pptx
- .txt/.md → 直接读取
- .png/.jpg/.jpeg → VLM 多模态分支

所有解析器输出统一的 UnifiedParserOutput。
"""

import logging
from pathlib import Path

from app.services.parser_interface import (
    DocumentParsingException,
    UnifiedParserOutput,
)

logger = logging.getLogger(__name__)


def parse_document(file_path: str, kind: str) -> str:
    """向后兼容的接口：返回纯文本字符串。

    内部调用 parse_document_full() 获取结构化输出，然后提取 markdown_content。
    新代码应直接使用 parse_document_full()。
    """
    output = parse_document_full(file_path, kind)
    return output.markdown_content


def parse_document_full(file_path: str, kind: str) -> UnifiedParserOutput:
    """完整解析接口：返回结构化的 UnifiedParserOutput。

    Args:
        file_path: 文件的物理路径
        kind: 材料类型 ("pdf", "ppt", "text_note", "image")

    Returns:
        UnifiedParserOutput

    Raises:
        ValueError: 不支持的材料类型
        DocumentParsingException: 解析失败
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    storage_root = _get_storage_root()

    # PDF: 双轨制路由
    if kind == "pdf":
        return _parse_pdf(path, storage_root)

    # Word/Excel/HTML
    if kind == "text_note" and ext in (".docx", ".doc", ".xlsx", ".html", ".htm"):
        # DOCX/DOC: 优先 MinerU V4（原生支持），降级到 MarkItDown
        if ext in (".docx", ".doc"):
            return _parse_office(path, storage_root, ext, fallback="markitdown")
        # XLSX/HTML: MarkItDown（MinerU V4 不支持 xlsx）
        return _parse_markitdown(path, ext)

    # 纯文本/Markdown: 直接读取
    if kind == "text_note":
        return _parse_text(path)

    # PPT: 优先 MinerU V4（原生支持），降级到 python-pptx
    if kind == "ppt":
        return _parse_office(path, storage_root, ext, fallback="pptx")

    # 图片: VLM 多模态分支
    if kind == "image":
        return _parse_image(path, storage_root)

    raise ValueError(f"Unsupported material kind: {kind}")


def _parse_pdf(path: Path, storage_root: Path) -> UnifiedParserOutput:
    """PDF 解析路由：MinerU (V4→V1) → Docling → pdfplumber。

    MinerU 公式识别质量远超 Docling，作为首选。
    仅当 MinerU 不可用（无 token / 全失败）时降级到本地解析器。
    """
    from app.core.config import settings

    # 优先走 MinerU（内部已包含 V4→V1 自动降级）
    has_mineru_token = bool(settings.mineru_api_token)
    if has_mineru_token:
        try:
            from app.services.mineru import MinerUParser
            parser = MinerUParser()
            return parser.parse(path, storage_root)
        except Exception as e:
            logger.warning("MinerU failed for %s: %s, falling back to local parsers", path.name, e)

    # 降级到 Docling（本地，电子版 PDF 效果好）
    try:
        from app.services.parsing_docling import DoclingParser
        parser = DoclingParser()
        return parser.parse(path, storage_root)
    except Exception as e:
        logger.warning("Docling failed for %s: %s, falling back to pdfplumber", path.name, e)

    # 最终兜底：pdfplumber
    return _parse_pdf_pdfplumber(path)


def _parse_pdf_pdfplumber(path: Path) -> UnifiedParserOutput:
    """pdfplumber 兜底解析（纯文字提取）。"""
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                pages.append(f"<!-- PAGE_START {i} -->\n{text}\n<!-- PAGE_END {i} -->")

    if not pages:
        raise DocumentParsingException(path, "pdfplumber extracted no text")

    return UnifiedParserOutput(
        raw_input_type="DIGITAL_PDF",
        markdown_content="\n\n".join(pages),
        page_count=len(pages),
        parser_name="pdfplumber",
    )


def _parse_office(
    path: Path, storage_root: Path, ext: str, fallback: str = "markitdown"
) -> UnifiedParserOutput:
    """Office 文档解析：优先 MinerU V4（原生支持 DOC/DOCX/PPT/PPTX），降级到本地解析器。"""
    from app.core.config import settings

    input_type_map = {
        ".docx": "WORD", ".doc": "WORD",
        ".pptx": "PPT", ".ppt": "PPT",
    }

    if settings.mineru_api_token:
        try:
            from app.services.mineru import MinerUParser
            parser = MinerUParser()
            output = parser.parse(path, storage_root)
            output.raw_input_type = input_type_map.get(ext, "TEXT")
            return output
        except Exception as e:
            logger.warning("MinerU failed for %s: %s, falling back to local parser", path.name, e)

    if fallback == "pptx":
        return _parse_pptx(path)
    return _parse_markitdown(path, ext)


def _parse_markitdown(path: Path, ext: str) -> UnifiedParserOutput:
    """MarkItDown 轻量解析（Word/Excel/HTML）。"""
    from markitdown import MarkItDown

    input_type_map = {
        ".docx": "WORD",
        ".xlsx": "EXCEL",
        ".html": "HTML",
        ".htm": "HTML",
    }

    try:
        md = MarkItDown()
        result = md.convert(str(path))
        text = result.text_content or ""
    except Exception as e:
        raise DocumentParsingException(path, "MarkItDown failed", e)

    if not text.strip():
        raise DocumentParsingException(path, "MarkItDown returned empty content")

    return UnifiedParserOutput(
        raw_input_type=input_type_map.get(ext, "TEXT"),
        markdown_content=text,
        parser_name="MarkItDown",
    )


def _parse_text(path: Path) -> UnifiedParserOutput:
    """纯文本/Markdown 直接读取。"""
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="gbk", errors="replace")
    except Exception as e:
        raise DocumentParsingException(path, "Failed to read text file", e)

    return UnifiedParserOutput(
        raw_input_type="TEXT",
        markdown_content=text,
        parser_name="text_reader",
    )


def _parse_pptx(path: Path) -> UnifiedParserOutput:
    """python-pptx 解析 PPT。"""
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise DocumentParsingException(path, "Failed to open PPTX", e)

    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(
                f"<!-- PAGE_START {i} -->\n"
                f"## Slide {i}\n" + "\n".join(texts) +
                f"\n<!-- PAGE_END {i} -->"
            )

    if not slides:
        raise DocumentParsingException(path, "PPTX contains no text content")

    return UnifiedParserOutput(
        raw_input_type="PPT",
        markdown_content="\n\n".join(slides),
        page_count=len(slides),
        parser_name="python-pptx",
    )


def _parse_image(path: Path, storage_root: Path) -> UnifiedParserOutput:
    """图片解析：MinerU V4/V1（原生支持 PNG/JPG），降级返回空内容提示。

    DeepSeek 当前无可用的视觉模型（v4-flash/v4-pro 均为纯文本），
    因此不走 VLM 分支，直接用 MinerU 的 OCR 能力处理图片。
    """
    from app.core.config import settings

    if settings.mineru_api_token:
        try:
            from app.services.mineru import MinerUParser
            parser = MinerUParser()
            output = parser.parse(path, storage_root)
            output.raw_input_type = "USER_IMAGE"
            return output
        except Exception as e:
            logger.warning("MinerU image parsing failed for %s: %s", path.name, e)

    # MinerU 不可用时，返回提示信息而非静默失败
    raise DocumentParsingException(
        path,
        "图片解析需要 MinerU API（当前未配置或调用失败）。"
        "请配置 MINERU_API_TOKEN 后重试。",
    )


def _get_storage_root() -> Path:
    """获取图片存储根目录。"""
    from app.core.config import settings
    storage_root = Path(settings.upload_root) / "storage" / "images"
    storage_root.mkdir(parents=True, exist_ok=True)
    return storage_root
